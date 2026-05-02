"""
Script 20 — Build presentation slides as a .pptx file.

Generates: results/presentation.pptx

Run from project root:
    source venv/bin/activate && python scripts/20_build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image as PILImage
import os
from pathlib import Path

RESULTS = 'results'

# ── Colour palette ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x0F, 0x1A)   # near-black navy
ACCENT    = RGBColor(0x4C, 0xAF, 0x50)   # green (our models)
ACCENT2   = RGBColor(0xFF, 0x7F, 0x0E)   # orange (graph baselines)
RED       = RGBColor(0xD6, 0x27, 0x28)   # red (sequence baselines)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGREY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def fit_dims(path, max_w, max_h):
    """Return (width, height) in EMU that fits path within max_w × max_h, preserving aspect ratio."""
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    return int(iw * scale), int(ih * scale)


def add_img(slide, path, left, top, width=None, height=None, max_w=None, max_h=None):
    if not os.path.exists(path):
        print(f'  WARNING: missing {path}')
        return
    if max_w or max_h:
        mw = max_w or Inches(13)
        mh = max_h or Inches(7)
        w, h = fit_dims(path, mw, mh)
        slide.shapes.add_picture(path, left, top, w, h)
    elif width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)


def add_rule(slide, left, top, width, color=ACCENT, thickness=Pt(2)):
    from pptx.util import Pt as Pt2
    line = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def title_slide(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Decorative top bar
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.08))

    add_text(slide,
             'Next-Chord Prediction via\nHeterogeneous Graph Neural Networks',
             Inches(1), Inches(1.5), Inches(11.3), Inches(2.2),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide,
             'McGill Billboard Dataset  ·  CES 941: Machine Learning on Graphs',
             Inches(1), Inches(3.7), Inches(11.3), Inches(0.6),
             size=20, color=LIGHTGREY, align=PP_ALIGN.CENTER)

    add_text(slide,
             'April 14 / 16, 2026',
             Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
             size=18, color=LIGHTGREY, align=PP_ALIGN.CENTER, italic=True)

    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.08))
    return slide


def section_header(prs, title, subtitle=''):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.08))

    # Large coloured box
    box = slide.shapes.add_shape(1, Inches(0), Inches(2.5), SLIDE_W, Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x1A)
    box.line.fill.background()

    add_text(slide, title, Inches(1), Inches(2.7), Inches(11.3), Inches(1.6),
             size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
                 size=20, color=LIGHTGREY, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.08))


def chord_box(slide, chord, left, top, box_w, box_h, bg_color, text_color=None):
    """Draw a rounded-corner chord pill."""
    from pptx.util import Pt as Pt2
    from pptx.enum.text import PP_ALIGN
    if text_color is None:
        text_color = WHITE
    shp = slide.shapes.add_shape(
        5,   # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        left, top, box_w, box_h,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg_color
    shp.line.fill.background()
    # corner rounding
    shp.adjustments[0] = 0.15
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = chord
    run.font.size = Pt2(13)
    run.font.bold = True
    run.font.color.rgb = text_color


def song_examples_slide(prs):
    """Slide 3 — Real song chord progressions as coloured chord boxes."""
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.06))

    add_text(slide, 'Chord Progressions You Already Know',
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=ACCENT)
    add_rule(slide, Inches(0.5), Inches(0.88), Inches(12.3), LIGHTGREY, Inches(0.015))
    add_text(slide, 'Each box is one chord — the dataset encodes exactly these kinds of sequences',
             Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.35),
             size=13, color=LIGHTGREY, italic=True)

    # Chord-type → box colour mapping (major=blue, minor=purple, 7th=orange, other=teal)
    MAJOR  = RGBColor(0x1F, 0x77, 0xB4)   # blue
    MINOR  = RGBColor(0x6A, 0x3D, 0x9A)   # purple
    DOM7   = RGBColor(0xFF, 0x7F, 0x0E)   # orange
    MAJ7   = RGBColor(0x2C, 0xA0, 0x2C)   # green
    SUS    = RGBColor(0x17, 0xBE, 0xCF)   # teal
    DIM    = RGBColor(0xD6, 0x27, 0x28)   # red

    # Each song: (title, artist, key, genre, note, [(chord_label, color), ...])
    songs = [
        (
            'Bohemian Rhapsody', 'Queen', 'B♭ major', 'Rock', 'Verse',
            [('B♭', MAJOR), ('Gm', MINOR), ('Cm', MINOR), ('F7', DOM7)],
        ),
        (
            'Piano Man', 'Billy Joel', 'C major', 'Pop/Rock', 'Chorus',
            [('C', MAJOR), ('F', MAJOR), ('Am', MINOR), ('G7', DOM7)],
        ),
        (
            'Africa', 'Toto', 'A major', 'Pop', 'Chorus',
            [('F♯m', MINOR), ('D', MAJOR), ('A', MAJOR), ('E', MAJOR)],
        ),
        (
            'Lose Yourself', 'Eminem', 'D minor', 'Hip-hop', 'Main riff',
            [('Dm', MINOR), ('B♭', MAJOR), ('C', MAJOR), ('A7', DOM7)],
        ),
        (
            'Something in the Orange', 'Zach Bryan', 'G major', 'Country', 'Verse',
            [('G', MAJOR), ('D', MAJOR), ('Em', MINOR), ('C', MAJOR)],
        ),
    ]

    BOX_W = Inches(1.35)
    BOX_H = Inches(0.46)
    ROW_H = Inches(1.22)
    START_Y = Inches(1.38)
    META_X  = Inches(0.5)
    CHORD_X = Inches(3.9)
    GAP     = Inches(0.10)

    for i, (title, artist, key, genre, note, chords) in enumerate(songs):
        y = START_Y + i * ROW_H

        # Song title + artist
        add_text(slide, title, META_X, y, Inches(3.2), Inches(0.38),
                 size=14, bold=True, color=WHITE)
        add_text(slide, f'{artist}  ·  {key}  ·  {genre}',
                 META_X, y + Inches(0.38), Inches(3.2), Inches(0.30),
                 size=11, color=LIGHTGREY)
        add_text(slide, note,
                 META_X, y + Inches(0.68), Inches(3.2), Inches(0.28),
                 size=10, color=YELLOW, italic=True)

        # Chord boxes
        cx = CHORD_X
        for chord, color in chords:
            chord_box(slide, chord, cx, y + Inches(0.08), BOX_W, BOX_H, color)
            cx += BOX_W + GAP

        # Row separator (not after last row)
        if i < len(songs) - 1:
            add_rule(slide, Inches(0.4), y + ROW_H - Inches(0.06),
                     Inches(12.5), RGBColor(0x33, 0x33, 0x44), Inches(0.01))

    # Legend
    LX = Inches(0.5)
    LY = Inches(7.08)
    legend = [
        ('Major', MAJOR), ('Minor', MINOR), ('Dom 7th', DOM7),
        ('Major 7th', MAJ7), ('Suspended', SUS), ('Diminished', DIM),
    ]
    add_text(slide, 'Chord quality:', LX, LY, Inches(1.3), Inches(0.35),
             size=10, color=LIGHTGREY)
    lx = LX + Inches(1.3)
    for label, col in legend:
        chord_box(slide, label, lx, LY, Inches(1.1), Inches(0.32), col)
        lx += Inches(1.15)

    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.06))
    return slide


def music_primer_slide(prs):
    """Slide 2 — Music Theory Primer: notes, intervals, chord types."""
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.06))

    # Title
    add_text(slide, 'Music Theory Primer',
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=ACCENT)
    add_rule(slide, Inches(0.5), Inches(0.88), Inches(12.3), LIGHTGREY, Inches(0.015))

    add_text(slide, 'Everything you need to follow this talk — no prior music knowledge required',
             Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.38),
             size=13, color=LIGHTGREY, italic=True)

    # ── Left column ──────────────────────────────────────────────────────────
    L = Inches(0.5)
    CW = Inches(5.7)

    add_text(slide, 'The 12 Notes  (one octave)',
             L, Inches(1.38), CW, Inches(0.38),
             size=15, bold=True, color=YELLOW)

    add_text(slide, 'C   C♯/D♭   D   D♯/E♭   E   F   F♯/G♭   G   G♯/A♭   A   A♯/B♭   B',
             L, Inches(1.78), CW, Inches(0.42),
             size=14, bold=True, color=WHITE)

    note_facts = [
        '• 12 equally-spaced pitches; the pattern repeats every octave',
        '• A semitone = one piano-key step (the smallest interval)',
        '• Sharp (♯) = one semitone up    Flat (♭) = one semitone down',
        '• C♯ and D♭ are the same pitch — "enharmonic equivalents"',
        '• After B the cycle restarts at C (one octave higher)',
    ]
    top = Inches(2.22)
    for fact in note_facts:
        add_text(slide, fact, L, top, CW, Inches(0.42), size=13, color=LIGHTGREY)
        top += Inches(0.40)

    add_text(slide, 'Key Concepts',
             L, Inches(4.35), CW, Inches(0.38),
             size=15, bold=True, color=YELLOW)

    concepts = [
        '• Chord — 3+ notes played simultaneously',
        '• Root note — the base pitch that names the chord (e.g. C in "C major")',
        '• Chord progression — a sequence of chords unfolding over time',
        '• Chord occurrence (occ) — one instance of a chord at a specific beat',
        '• Section — a musical segment (verse, chorus, bridge, intro …)',
    ]
    top = Inches(4.78)
    for c in concepts:
        add_text(slide, c, L, top, CW, Inches(0.42), size=13, color=LIGHTGREY)
        top += Inches(0.40)

    # ── Vertical divider ─────────────────────────────────────────────────────
    div = slide.shapes.add_shape(1, Inches(6.5), Inches(1.35), Inches(0.025), Inches(5.75))
    div.fill.solid(); div.fill.fore_color.rgb = LIGHTGREY
    div.line.fill.background()

    # ── Right column ─────────────────────────────────────────────────────────
    R = Inches(6.75)
    RW = Inches(6.3)

    add_text(slide, 'Chord Types in This Dataset',
             R, Inches(1.38), RW, Inches(0.38),
             size=15, bold=True, color=YELLOW)

    add_text(slide, '144 unique chord types  =  12 root notes  ×  ~12 quality variants',
             R, Inches(1.78), RW, Inches(0.38),
             size=12, color=LIGHTGREY, italic=True)

    # Table: (label, notes, description)
    chord_rows = [
        ('Major          (C)',       'C – E – G',         'bright, stable — the "happy" default'),
        ('Minor          (Cm)',      'C – E♭ – G',        'darker, "sadder" feel'),
        ('Dominant 7th   (C7)',      'C – E – G – B♭',    'bluesy tension, wants to resolve'),
        ('Major 7th      (Cmaj7)',   'C – E – G – B',     'smooth, jazzy, open'),
        ('Minor 7th      (Cm7)',     'C – E♭ – G – B♭',  'mellow, soulful'),
        ('Diminished     (Cdim)',    'C – E♭ – G♭',       'tense, unstable'),
        ('Augmented      (Caug)',    'C – E – G♯',        'dreamy, unresolved'),
        ('Suspended 4th  (Csus4)',   'C – F – G',         'open, no major/minor quality'),
        ('Suspended 2nd  (Csus2)',   'C – D – G',         'airy, spacious'),
        ('Power chord    (C5)',      'C – G',             'root + 5th only; raw rock sound'),
        ('Add9           (Cadd9)',   'C – E – G – D',     'colour tone added above the triad'),
        ('No-chord       (N)',       '—',                 'silence or unidentifiable harmony'),
    ]

    # Column headers
    col_label = R
    col_notes = R + Inches(2.25)
    col_desc  = R + Inches(3.75)

    add_text(slide, 'Name',  col_label, Inches(2.18), Inches(2.2),  Inches(0.32), size=11, bold=True, color=ACCENT)
    add_text(slide, 'Notes', col_notes, Inches(2.18), Inches(1.45), Inches(0.32), size=11, bold=True, color=ACCENT)
    add_text(slide, 'Sound', col_desc,  Inches(2.18), Inches(2.5),  Inches(0.32), size=11, bold=True, color=ACCENT)
    add_rule(slide, R, Inches(2.52), RW, ACCENT, Inches(0.01))

    top = Inches(2.57)
    row_h = Inches(0.375)
    for i, (label, notes, desc) in enumerate(chord_rows):
        row_color = WHITE if i % 2 == 0 else LIGHTGREY
        add_text(slide, label, col_label, top, Inches(2.2),  Inches(row_h), size=11, color=row_color)
        add_text(slide, notes, col_notes, top, Inches(1.45), Inches(row_h), size=11, color=row_color)
        add_text(slide, desc,  col_desc,  top, Inches(2.5),  Inches(row_h), size=11, color=row_color)
        top += row_h

    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.06))
    return slide


def content_slide(prs, title, bullets=None, img_path=None,
                  img_left=None, img_top=None, img_width=None, img_height=None,
                  bullet_left=Inches(0.5), bullet_width=Inches(5.8),
                  note=None):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.06))

    # Title
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=ACCENT)
    # Thin separator under title
    add_rule(slide, Inches(0.5), Inches(0.88), Inches(12.3), LIGHTGREY, Inches(0.015))

    # Bullets
    if bullets:
        top = Inches(1.05)
        for bullet in bullets:
            if bullet == '':          # use as a vertical spacer, no shape added
                top += Inches(0.2)
                continue
            indent = bullet.startswith('  ')
            txt    = bullet.lstrip()
            prefix = '  • ' if not indent else '      – '
            tb = slide.shapes.add_textbox(bullet_left, top, bullet_width, Inches(0.55))
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = prefix + txt
            run.font.size  = Pt(18 if not indent else 15)
            run.font.bold  = False
            run.font.color.rgb = WHITE if not indent else LIGHTGREY
            top += Inches(0.52 if not indent else 0.42)

    # Image
    if img_path:
        add_img(slide, img_path,
                img_left  or Inches(6.5),
                img_top   or Inches(0.95),
                img_width, img_height)

    # Footer note
    if note:
        add_text(slide, note,
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                 size=12, color=LIGHTGREY, italic=True)

    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.06))
    return slide


def full_img_slide(prs, title, img_path, note=None):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.06))
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=ACCENT)
    add_rule(slide, Inches(0.5), Inches(0.88), Inches(12.3), LIGHTGREY, Inches(0.015))
    # Fit within the content area: max 12.4" wide, 6.2" tall, centred horizontally
    max_w, max_h = Inches(12.4), Inches(6.2)
    w, h = fit_dims(img_path, max_w, max_h)
    left = (SLIDE_W - w) // 2
    add_img(slide, img_path, left, Inches(0.97), max_w=max_w, max_h=max_h)
    if note:
        add_text(slide, note,
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                 size=12, color=LIGHTGREY, italic=True)
    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.06))
    return slide


def two_img_slide(prs, title, left_path, right_path,
                  left_label='', right_label='', note=None):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rule(slide, Inches(0), Inches(0), SLIDE_W, ACCENT, Inches(0.06))
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=ACCENT)
    add_rule(slide, Inches(0.5), Inches(0.88), Inches(12.3), LIGHTGREY, Inches(0.015))
    # Each image gets roughly half the slide width
    half_max_w, half_max_h = Inches(6.2), Inches(5.7)
    label_top = Inches(0.93)
    img_top   = Inches(1.35)
    if left_label:
        add_text(slide, left_label, Inches(0.2), label_top, Inches(6.2), Inches(0.38),
                 size=13, color=LIGHTGREY, align=PP_ALIGN.CENTER)
    if right_label:
        add_text(slide, right_label, Inches(6.8), label_top, Inches(6.2), Inches(0.38),
                 size=13, color=LIGHTGREY, align=PP_ALIGN.CENTER)
    add_img(slide, left_path,  Inches(0.2),  img_top, max_w=half_max_w, max_h=half_max_h)
    add_img(slide, right_path, Inches(6.85), img_top, max_w=half_max_w, max_h=half_max_h)
    if note:
        add_text(slide, note,
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                 size=12, color=LIGHTGREY, italic=True)
    add_rule(slide, Inches(0), Inches(7.4), SLIDE_W, ACCENT, Inches(0.06))
    return slide


# ── Build deck ────────────────────────────────────────────────────────────────

def build(prs):

    # 1 — Title
    title_slide(prs)

    # 2 — Music Theory Primer
    music_primer_slide(prs)

    # 3 — Song examples (chord progression boxes)
    song_examples_slide(prs)

    # 4 — Problem & Dataset
    content_slide(prs,
        title='Problem & Dataset',
        bullets=[
            'Task: given the current chord occurrence, predict the next chord',
            'McGill Billboard corpus: 890 pop/rock songs (~100K chord occurrences)',
            '144 unique chord types + N (no-chord) = 145-class prediction',
            'Why graph?',
            '  Chord progressions have rich hierarchical structure: occurrence → section → song',
            '  Flat sequence models discard this structure entirely',
            '  A heterogeneous graph can model all relationships simultaneously',
        ],
        img_path=f'{RESULTS}/genre_chord_heatmap.png',
        img_left=Inches(7.0), img_top=Inches(1.0),
        img_width=Inches(5.9),
        note='Genre distribution: rock 52%, pop 19%, soul/R&B 7%, country 4%, other 18%',
    )

    # 3a — Graph schema
    full_img_slide(prs,
        title='Graph Construction — Schema',
        img_path=f'{RESULTS}/graph_schema.png',
        note='5 node types · 10 directed edge types · all edges are within a single song graph',
    )

    # 3b — Graph instance
    full_img_slide(prs,
        title='Graph Construction — Real Song Example',
        img_path=f'{RESULTS}/graph_instance.png',
        note='occ nodes form a timeline; chord nodes are shared within the song; note/scale_deg nodes are global music-theory anchors',
    )

    # 4 — Sequence baselines
    content_slide(prs,
        title='Sequence Baselines',
        bullets=[
            'Markov (bigram + section-conditioned):  23.5% top-1',
            'LSTM (2-layer, hidden=256):  45.2% top-1',
            'Transformer (3-layer causal, embed=128):  52.2% top-1',
            '',
            'Transformer beats LSTM by +7pp — attention helps on chord sequences',
            'Both still far behind our GNN — the graph adds something real',
        ],
        img_path=f'{RESULTS}/model_comparison_final.png',
        img_left=Inches(6.4), img_top=Inches(1.0),
        img_width=Inches(6.8),
        note='All models trained on same 623/133/134 train/val/test split',
    )

    # 5 — The leakage slide
    content_slide(prs,
        title='The GNN Gets 99.8% — Something Is Wrong',
        bullets=[
            'Initial full GNN (with prev edges): 99.8% top-1',
            'LSTM: 45.2%.  Markov: 23.5%.  This gap is suspicious.',
            '',
            'Root cause: the prev edge creates a 2-hop leakage path',
            '  Layer 1:  chord[i+1] →(inst_rev)→ occ[i+1]',
            '    chord identity flows into occ[i+1] embedding',
            '  Layer 2:  occ[i+1] →(prev)→ occ[i]',
            '    future chord identity reaches the node we are predicting from',
            '',
            'Ablation confirms it: remove instance_of → accuracy collapses to 7.5%',
            '  Without chord identity in occ embeddings, the leak carries nothing',
        ],
        bullet_width=Inches(8.5),
        note='prev edge was added to allow bidirectional context — it accidentally enabled data leakage',
    )

    # 6 — Causal fix
    content_slide(prs,
        title='Fix: Causal Model (No prev Edges)',
        bullets=[
            'Remove prev edges entirely → no future information can flow back',
            'Inject current chord features directly into occ input vector',
            '  occ input: [timing(19-d) | chord_features(25-d)] = 44-d',
            '  This gives the model the present chord without leaking the future',
            '3 SAGEConv layers, hidden_dim=128, dropout=0.3',
            '',
            'Result:  60.7% top-1  ·  89.0% top-5  ·  CE 1.44',
            'Legitimate learning — beats Transformer by +8.5pp',
        ],
        img_path=f'{RESULTS}/causal_gnn_training_curves.png',
        img_left=Inches(7.0), img_top=Inches(1.0),
        img_width=Inches(6.0),
        note='Checkpoint: results/causal_gnn_best.pt',
    )

    # 7 — Full results table
    full_img_slide(prs,
        title='All Results — Test Set',
        img_path=f'{RESULTS}/results_table.png',
        note='Green = our models. Orange = graph baselines. Red = sequence baselines.',
    )

    # 8 — Design ladder
    full_img_slide(prs,
        title='Incremental Gain from Each Design Decision',
        img_path=f'{RESULTS}/design_ladder.png',
        note='+37.2pp total gain from Markov → Het. GNN.  Each architectural choice contributes.',
    )

    # 9 — Ablation
    content_slide(prs,
        title='Ablation: What Actually Matters?',
        bullets=[
            'Remove seq edges (next/prev):  25.4%  — sequencing is essential',
            'Remove instance_of:  7.5%  — chord identity is the dominant signal',
            '',
            'Homo GNN vs Het. GNN (key comparison):',
            '  Homo GNN (flat conv, all edges same type):  55.7%',
            '  Het. GNN v2 (type-specialised SAGEConv):  60.7%',
            '  +5pp from heterogeneous type specialisation',
            '',
            'HGT (Hu et al. 2020) vs Het. GNN:',
            '  HGT:  55.5%  ≈  HomoGNN:  55.7%',
            '  Graph structure matters more than the attention mechanism',
        ],
        img_path=f'{RESULTS}/ablation_final.png',
        img_left=Inches(7.0), img_top=Inches(1.0),
        img_width=Inches(6.1),
        bullet_width=Inches(6.3),
    )

    # 10 — GAT attention
    content_slide(prs,
        title='GAT Attention — What Does the Model Focus On?',
        bullets=[
            'GAT (4 heads) replaces SAGEConv — learns per-edge attention weights',
            'Attention is normalised per destination node across all incoming edge types',
            '',
            'inst_rev (chord → occ):  ≈ 1.0  — chord identity is the dominant signal',
            'next (prev occ → occ):   ≈ 1.0  — sequential context equally important',
            'next_section:            ≈ 1.0  — section transitions also attended to',
            'in_section (occ → sec):  ≈ 0.13 — section context used but secondary',
            'instance_of:             ≈ 0.13 — low weight (inst_rev carries same info)',
            '',
            'GAT top-1: 60.3%  vs  SAGEConv: 60.7% — marginal difference',
            '  Attention learns meaningful weights but does not improve accuracy',
            '  Confirms: graph structure matters more than weighting mechanism',
        ],
        img_path=f'{RESULTS}/gat_attention_by_relation.png',
        img_left=Inches(6.9), img_top=Inches(1.0),
        img_width=Inches(6.2),
        bullet_width=Inches(6.2),
        note='Attention extracted from layer 0 across all edge types simultaneously',
    )

    # 11 — Section accuracy
    content_slide(prs,
        title='Section-Level Accuracy',
        bullets=[
            'Instrumental / intro / outro:  73–74%  (repetitive, predictable)',
            'Verse / chorus:  58–59%  (moderate harmonic variety)',
            'Bridge:  51.7%  (most harmonically diverse — hardest to predict)',
            'Trans / interlude / solo:  66–72%  (shorter, often formulaic)',
            '',
            'Model captures musical structure: easier sections → higher accuracy',
        ],
        img_path=f'{RESULTS}/causal_gnn_section_accuracy.png',
        img_left=Inches(6.8), img_top=Inches(1.0),
        img_width=Inches(6.3),
        note='Section type decoded from sec node features via one-hot argmax',
    )

    # 12 — Genre probe results
    content_slide(prs,
        title='Genre Classification',
        bullets=[
            'Transformer (linear probe, frozen):  23.9%',
            'Homo GNN (linear probe, frozen):  23.9%  — near-random baseline ≈ 18%',
            '  → chord-prediction models do not spontaneously encode genre',
            '',
            'Het. GNN + genre head (joint training):  35.8%',
            '  → genre supervision adds +12pp over probe',
            '',
            'HGT + genre head (hierarchical sec pooling):  50.7%',
            '  → pooling over sections (not raw occ) gives richer song-level repr.',
            '  → sec nodes absorb occ harmonic context through in_section edges',
        ],
        img_path=f'{RESULTS}/genre_probe_comparison.png',
        img_left=Inches(7.0), img_top=Inches(1.05),
        img_width=Inches(6.1),
        bullet_width=Inches(6.3),
        note='Genre accuracy: 9-class classification on 134 test songs',
    )

    # 13 — Conclusion
    content_slide(prs,
        title='Conclusion',
        bullets=[
            'Heterogeneous graph > sequence models',
            '  60.7% vs 52.2% Transformer  (+8.5pp)',
            'Type specialisation matters',
            '  Het. GNN 60.7% vs Homo GNN 55.7%  (+5pp)',
            'Graph structure > attention mechanism',
            '  HGT 55.5% ≈ Homo GNN 55.7% — how you type edges beats how you weight them',
            'Leakage is detectable and diagnosable',
            '  99.8% → ablation → causal fix → 60.7% legitimate accuracy',
            'Section structure enables genre understanding',
            '  Hierarchical sec pooling: 50.7% genre vs 23.9% probe baseline',
            '',
            'Future work: contrastive learning for genre-separable embeddings',
        ],
        bullet_width=Inches(8.0),
        note='Code + results: github.com/[repo] | Checkpoints: results/*.pt',
    )

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == '__main__':
    Path(RESULTS).mkdir(exist_ok=True)
    prs  = new_prs()
    build(prs)
    out  = f'{RESULTS}/presentation.pptx'
    prs.save(out)
    print(f'Saved {out}  ({len(prs.slides)} slides)')
